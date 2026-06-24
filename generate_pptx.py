#!/usr/bin/env python3
"""Generate a professional pitch deck for 智慧校园管理系统."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree
import os

# Color palette
PRIMARY_BLUE   = RGBColor(0x25, 0x63, 0xeb)
CYAN_ACCENT    = RGBColor(0x06, 0xb, 0xd6)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
BLACK          = RGBColor(0x11, 0x11, 0x11)
DARK_GRAY      = RGBColor(0x37, 0x41, 0x51)
MED_GRAY       = RGBColor(0x6B, 0x72, 0x80)
LIGHT_GRAY     = RGBColor(0xF3, 0xF4, 0xF6)
VERY_LIGHT_BLUE = RGBColor(0xEF, 0xF2, 0xFE)
SUBTLE_CYAN    = RGBColor(0xE0, 0xF2, 0xFE)

TITLE_FONT      = "Microsoft YaHei"
BODY_FONT       = "Microsoft YaHei"
NUM_FONT        = "Arial"

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


def add_background(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_accent_bar(slide, left, top, width, height, color=PRIMARY_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=BODY_FONT, line_spacing=1.3, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return txBox


def tb(slide, left, top, width, height, text, font_size=18,
       font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
       font_name=BODY_FONT, line_spacing=1.3):
    """Convenience wrapper using keyword args only."""
    return add_textbox(slide, left, top, width, height, text, font_size,
                       font_color, bold, alignment, font_name, line_spacing)


def add_circle_number(slide, left, top, number, color=PRIMARY_BLUE, size=Inches(0.55)):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = NUM_FONT
    p.alignment = PP_ALIGN.CENTER
    return circle


def add_rounded_card(slide, left, top, width, height, bg_color=VERY_LIGHT_BLUE,
                     border_color=None, border_width_pt=0):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(border_width_pt)
    else:
        card.line.fill.background()
    return card


# ════════════════════════════════════════════════════════════════
# SLIDE 1 - Title Slide
# ════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide1, PRIMARY_BLUE)
add_shape_bg(slide1, Inches(0), Inches(0), Inches(4.5), SLIDE_HEIGHT, CYAN_ACCENT)

add_textbox(slide1, Inches(1.2), Inches(1.8), Inches(10), Inches(1.2),
            "智慧校园管理系统", 48, WHITE, True, PP_ALIGN.CENTER, font_name=TITLE_FONT)

add_accent_bar(slide1, Inches(5.2), Inches(3.1), Inches(3), Inches(0.06), CYAN_ACCENT)

add_textbox(slide1, Inches(1.2), Inches(3.3), Inches(10), Inches(0.8),
            "一站式数字化校园解决方案", 28, RGBColor(0xBB, 0xDE, 0xFB), False, PP_ALIGN.CENTER, font_name=BODY_FONT)

add_textbox(slide1, Inches(1.2), Inches(5.2), Inches(10), Inches(0.6),
            "参赛团队", 20, RGBColor(0x93, 0xC5, 0xFD), False, PP_ALIGN.CENTER, font_name=BODY_FONT)

add_accent_bar(slide1, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)


# ════════════════════════════════════════════════════════════════
# SLIDE 2 - Industry Pain Points
# ════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, WHITE)
add_shape_bg(slide2, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

add_textbox(slide2, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "行业痛点", 14, CYAN_ACCENT, True, font_name=BODY_FONT)
add_textbox(slide2, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "校园管理面临的四大挑战", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide2, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), PRIMARY_BLUE)

pain_data = [
    ("数据孤岛", "校园管理信息化程度低，\n数据孤岛严重", Inches(0.8), Inches(2.2)),
    ("效率低下", "师生日常办事流程繁琐，\n效率低下", Inches(7.0), Inches(2.2)),
    ("缺乏统一", "各部门系统独立，\n缺乏统一管理平台", Inches(0.8), Inches(4.3)),
    ("难以支撑", "传统管理方式难以支撑\n现代化教学需求", Inches(7.0), Inches(4.3)),
]

for i, (label, desc, px, py) in enumerate(pain_data):
    card = add_rounded_card(slide2, px, py, Inches(5.8), Inches(1.7), VERY_LIGHT_BLUE)
    add_circle_number(slide2, px + Inches(0.2), py + Inches(0.2), i + 1)
    add_textbox(slide2, px + Inches(0.85), py + Inches(0.2), Inches(4.5), Inches(0.45),
                label, 20, PRIMARY_BLUE, True, font_name=BODY_FONT)
    add_textbox(slide2, px + Inches(0.85), py + Inches(0.75), Inches(4.5), Inches(0.8),
                desc, 16, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 3 - Our Solution
# ════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, WHITE)
add_shape_bg(slide3, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)

add_textbox(slide3, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "解决方案", 14, PRIMARY_BLUE, True, font_name=BODY_FONT)
add_textbox(slide3, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "我们的破局之道", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide3, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), CYAN_ACCENT)

sol_data = [
    ("一体化平台", "打通数据壁垒，实现信息互联", Inches(0.8), Inches(2.2), CYAN_ACCENT),
    ("全场景覆盖", "教务/宿舍/餐饮/图书馆核心场景", Inches(7.0), Inches(2.2), PRIMARY_BLUE),
    ("数据看板", "实时数据可视化，辅助科学决策", Inches(0.8), Inches(4.3), CYAN_ACCENT),
    ("多端适配", "响应式设计，PC/移动端无缝访问", Inches(7.0), Inches(4.3), PRIMARY_BLUE),
]

for label, desc, px, py, color in sol_data:
    card = add_rounded_card(slide3, px, py, Inches(5.8), Inches(1.7), SUBTLE_CYAN)
    # Icon circle
    icon = slide3.shapes.add_shape(MSO_SHAPE.OVAL, px + Inches(0.2), py + Inches(0.2), Inches(0.5), Inches(0.5))
    icon.fill.solid()
    icon.fill.fore_color.rgb = color
    icon.line.fill.background()
    tf = icon.text_frame
    p = tf.paragraphs[0]
    p.text = "+"
    p.font.size = Pt(22)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.CENTER
    add_textbox(slide3, px + Inches(0.85), py + Inches(0.2), Inches(4.5), Inches(0.45),
                label, 20, color, True, font_name=BODY_FONT)
    add_textbox(slide3, px + Inches(0.85), py + Inches(0.75), Inches(4.5), Inches(0.8),
                desc, 16, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 4 - Core Features
# ════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide4, WHITE)
add_shape_bg(slide4, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

add_textbox(slide4, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "核心功能", 14, CYAN_ACCENT, True, font_name=BODY_FONT)
add_textbox(slide4, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "七大核心模块，覆盖校园全场景", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide4, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), PRIMARY_BLUE)

features = [
    ("数据看板", "实时运营概览\n可视化图表"),
    ("教务管理", "课程、成绩、\n教学安排"),
    ("课程表", "智能排课\n周历视图"),
    ("学生管理", "学籍、信息、\n状态追踪"),
    ("图书馆", "馆藏检索、\n借阅管理"),
    ("餐饮服务", "菜单、\n营养分析"),
    ("宿舍管理", "分配、报修、\n入住统计"),
]

cols = 4
rows = 2
card_w = Inches(2.8)
card_h = Inches(1.5)
gap_x = Inches(0.35)
gap_y = Inches(0.35)
start_x = Inches(0.8)
start_y = Inches(2.15)

for idx, (title, desc) in enumerate(features):
    col = idx % cols
    row = idx // cols
    cx = start_x + col * (card_w + gap_x)
    cy = start_y + row * (card_h + gap_y)

    bg = VERY_LIGHT_BLUE if row == 0 else SUBTLE_CYAN
    card = add_rounded_card(slide4, cx, cy, card_w, card_h, bg, RGBColor(0xD1, 0xD5, 0xDB), 1)

    # Small colored dot
    dot = slide4.shapes.add_shape(MSO_SHAPE.OVAL, cx + Inches(0.15), cy + Inches(0.12), Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = PRIMARY_BLUE
    dot.line.fill.background()

    add_textbox(slide4, cx + Inches(0.1), cy + Inches(0.45), card_w - Inches(0.2), Inches(0.35),
                title, 15, BLACK, True, PP_ALIGN.CENTER, font_name=BODY_FONT)
    add_textbox(slide4, cx + Inches(0.1), cy + Inches(0.85), card_w - Inches(0.2), Inches(0.55),
                desc, 12, MED_GRAY, False, PP_ALIGN.CENTER, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 5 - Tech Architecture
# ════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide5, WHITE)
add_shape_bg(slide5, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)

add_textbox(slide5, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "技术架构", 14, PRIMARY_BLUE, True, font_name=BODY_FONT)
add_textbox(slide5, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "现代化前端技术栈", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide5, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), CYAN_ACCENT)

tech_items = [
    ("React 18 + Vite 5", "高性能前端框架", PRIMARY_BLUE),
    ("Tailwind CSS 3", "原子化样式方案", CYAN_ACCENT),
    ("Recharts", "可视化图表库", PRIMARY_BLUE),
    ("React Router DOM", "声明式路由管理", CYAN_ACCENT),
    ("Lucide React", "现代化图标库", PRIMARY_BLUE),
]

positions = [
    (Inches(0.8), Inches(2.2)),
    (Inches(7.0), Inches(2.2)),
    (Inches(0.8), Inches(4.3)),
    (Inches(4.2), Inches(4.3)),
    (Inches(7.0), Inches(4.3)),
]

for i, (tech, desc, color) in enumerate(tech_items):
    px, py = positions[i]
    card = add_rounded_card(slide5, px, py, Inches(2.6), Inches(1.7), WHITE, color, 2)
    add_accent_bar(slide5, px, py, Inches(2.6), Inches(0.06), color)
    add_textbox(slide5, px + Inches(0.15), py + Inches(0.2), Inches(2.3), Inches(0.45),
                tech, 18, BLACK, True, PP_ALIGN.CENTER, font_name=BODY_FONT)
    add_textbox(slide5, px + Inches(0.15), py + Inches(0.7), Inches(2.3), Inches(0.8),
                desc, 14, MED_GRAY, False, PP_ALIGN.CENTER, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 6 - Product Highlights
# ════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide6, WHITE)
add_shape_bg(slide6, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

add_textbox(slide6, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "产品亮点", 14, CYAN_ACCENT, True, font_name=BODY_FONT)
add_textbox(slide6, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "为什么选择我们", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide6, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), PRIMARY_BLUE)

highlights = [
    ("7大核心模块", "覆盖校园全场景，从教务到后勤一站式解决", PRIMARY_BLUE),
    ("响应式设计", "自适应PC、平板、手机，随时随地访问", CYAN_ACCENT),
    ("组件化架构", "高内聚低耦合，易于扩展和维护升级", PRIMARY_BLUE),
    ("实时可视化", "数据驱动决策，运营状态一目了然", CYAN_ACCENT),
]

for i, (title, desc, color) in enumerate(highlights):
    y = Inches(2.2) + i * Inches(1.2)
    add_accent_bar(slide6, Inches(0.8), y, Inches(0.08), Inches(0.9), color)
    add_shape_bg(slide6, Inches(0.95), y, Inches(11.2), Inches(0.9), LIGHT_GRAY)
    add_textbox(slide6, Inches(1.3), y + Inches(0.05), Inches(10), Inches(0.45),
                title, 22, BLACK, True, font_name=BODY_FONT)
    add_textbox(slide6, Inches(1.3), y + Inches(0.48), Inches(10), Inches(0.4),
                desc, 16, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 7 - Market Analysis
# ════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide7, WHITE)
add_shape_bg(slide7, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)

add_textbox(slide7, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "市场前景", 14, PRIMARY_BLUE, True, font_name=BODY_FONT)
add_textbox(slide7, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "千亿级蓝海市场", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide7, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), CYAN_ACCENT)

market_items = [
    ("千亿市场", "中国智慧校园市场规模超千亿", Inches(0.8), Inches(2.2)),
    ("政策驱动", "教育数字化转型全面加速", Inches(7.0), Inches(2.2)),
    ("融合需求", "线上线下混合教学成为常态", Inches(0.8), Inches(4.3)),
    ("目标客户", "高校、中学、职业院校全覆盖", Inches(7.0), Inches(4.3)),
]

for i, (label, desc, px, py) in enumerate(market_items):
    card = add_rounded_card(slide7, px, py, Inches(5.8), Inches(1.7), VERY_LIGHT_BLUE)
    # Number badge
    num = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, px + Inches(0.15), py + Inches(0.15), Inches(0.6), Inches(0.3))
    num.fill.solid()
    num.fill.fore_color.rgb = PRIMARY_BLUE
    num.line.fill.background()
    tf = num.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(slide7, px + Inches(0.85), py + Inches(0.15), Inches(4.5), Inches(0.4),
                label, 20, PRIMARY_BLUE, True, font_name=BODY_FONT)
    add_textbox(slide7, px + Inches(0.85), py + Inches(0.65), Inches(4.5), Inches(0.8),
                desc, 15, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 8 - Competitive Edge
# ════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide8, WHITE)
add_shape_bg(slide8, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

add_textbox(slide8, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "竞争优势", 14, CYAN_ACCENT, True, font_name=BODY_FONT)
add_textbox(slide8, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "差异化竞争策略", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide8, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), PRIMARY_BLUE)

comp_data = [
    ("轻量化部署", "低成本快速上线，无需复杂硬件", PRIMARY_BLUE, Inches(0.8), Inches(2.2)),
    ("模块化设计", "按需组合功能模块，灵活适配", CYAN_ACCENT, Inches(7.0), Inches(2.2)),
    ("优秀体验", "简洁直观的操作界面，降低学习成本", PRIMARY_BLUE, Inches(0.8), Inches(4.3)),
    ("开源生态", "成熟技术栈，社区活跃，持续迭代", CYAN_ACCENT, Inches(7.0), Inches(4.3)),
]

for label, desc, color, px, py in comp_data:
    card = add_rounded_card(slide8, px, py, Inches(5.8), Inches(1.7), WHITE, color, 2)
    add_textbox(slide8, px + Inches(0.15), py + Inches(0.15), Inches(5.5), Inches(0.45),
                label, 20, BLACK, True, font_name=BODY_FONT)
    add_textbox(slide8, px + Inches(0.15), py + Inches(0.7), Inches(5.5), Inches(0.8),
                desc, 15, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)
    add_accent_bar(slide8, px, py + Inches(1.64), Inches(5.8), Inches(0.06), color)


# ════════════════════════════════════════════════════════════════
# SLIDE 9 - Business Model
# ════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide9, WHITE)
add_shape_bg(slide9, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)

add_textbox(slide9, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "商业模式", 14, PRIMARY_BLUE, True, font_name=BODY_FONT)
add_textbox(slide9, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "多元化的盈利模式", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide9, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), CYAN_ACCENT)

biz_data = [
    ("SaaS订阅", "按学校规模分级定价", "稳定现金流", PRIMARY_BLUE, Inches(0.8)),
    ("定制开发", "针对大型客户需求", "高客单价", CYAN_ACCENT, Inches(4.8)),
    ("增值服务", "AI推荐、数据分析", "高利润率", PRIMARY_BLUE, Inches(8.8)),
]

for label, desc, revenue, color, x_pos in biz_data:
    cw = Inches(3.4)
    card = add_rounded_card(slide9, x_pos, Inches(2.2), cw, Inches(3.2), WHITE, color, 2)
    add_shape_bg(slide9, x_pos, Inches(2.2), cw, Inches(0.6), color)
    add_textbox(slide9, x_pos, Inches(2.25), cw, Inches(0.5),
                revenue, 14, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide9, x_pos + Inches(0.15), Inches(2.95), cw - Inches(0.3), Inches(0.5),
                label, 22, BLACK, True, PP_ALIGN.CENTER, font_name=BODY_FONT)
    add_textbox(slide9, x_pos + Inches(0.15), Inches(3.5), cw - Inches(0.3), Inches(1.5),
                desc, 16, MED_GRAY, False, PP_ALIGN.CENTER, font_name=BODY_FONT, line_spacing=1.5)


# ════════════════════════════════════════════════════════════════
# SLIDE 10 - Roadmap
# ════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide10, WHITE)
add_shape_bg(slide10, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), PRIMARY_BLUE)

add_textbox(slide10, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "发展规划", 14, CYAN_ACCENT, True, font_name=BODY_FONT)
add_textbox(slide10, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "产品演进路线图", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide10, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), PRIMARY_BLUE)

# Timeline horizontal line
add_accent_bar(slide10, Inches(1.0), Inches(2.55), Inches(11.3), Inches(0.04), RGBColor(0xD1, 0xD5, 0xDB))

roadmap = [
    ("Phase 1", "核心功能完善", "已完成", PRIMARY_BLUE, True, Inches(0.8)),
    ("Phase 2", "移动端 APP / 小程序", "进行中", CYAN_ACCENT, False, Inches(4.0)),
    ("Phase 3", "AI 智能推荐与数据分析", "规划中", PRIMARY_BLUE, False, Inches(7.2)),
    ("Phase 4", "多校区、多租户支持", "愿景", CYAN_ACCENT, False, Inches(10.4)),
]

for phase, title, status, color, done, x_pos in roadmap:
    # Timeline dot
    dot = slide10.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(0.3), Inches(2.35), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    # Phase label above
    add_textbox(slide10, x_pos, Inches(2.0), Inches(3.0), Inches(0.4),
                phase, 14, MED_GRAY, True, PP_ALIGN.CENTER, font_name=BODY_FONT)

    # Title below
    add_textbox(slide10, x_pos, Inches(2.85), Inches(3.0), Inches(0.45),
                title, 18, BLACK, True, PP_ALIGN.CENTER, font_name=BODY_FONT)

    # Status badge
    badge = slide10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x_pos + Inches(0.5), Inches(3.4), Inches(2.0), Inches(0.35))
    badge.fill.solid()
    badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = BODY_FONT
    p.alignment = PP_ALIGN.CENTER

    # Done checkmark
    if done:
        check = slide10.shapes.add_shape(MSO_SHAPE.OVAL, x_pos + Inches(2.55), Inches(2.3), Inches(0.35), Inches(0.35))
        check.fill.solid()
        check.fill.fore_color.rgb = WHITE
        check.line.fill.background()
        tf2 = check.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = "OK"
        p2.font.size = Pt(9)
        p2.font.color.rgb = color
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════════
# SLIDE 11 - Team Introduction
# ════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide11, WHITE)
add_shape_bg(slide11, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)

add_textbox(slide11, Inches(0.8), Inches(0.4), Inches(11), Inches(0.5),
            "团队介绍", 14, PRIMARY_BLUE, True, font_name=BODY_FONT)
add_textbox(slide11, Inches(0.8), Inches(0.9), Inches(11), Inches(1.0),
            "我们的团队", 38, BLACK, True, font_name=TITLE_FONT)
add_accent_bar(slide11, Inches(0.8), Inches(1.7), Inches(2.5), Inches(0.05), CYAN_ACCENT)

team_items = [
    ("技术实力", "核心团队由计算机专业学生组成", Inches(0.8), Inches(2.2)),
    ("综合能力", "前端开发、产品设计、项目管理", Inches(7.0), Inches(2.2)),
    ("导师支持", "指导老师来自计算机学院", Inches(0.8), Inches(4.3)),
    ("创新精神", "以学生为中心，用技术改善校园生活", Inches(7.0), Inches(4.3)),
]

for label, desc, px, py in team_items:
    card = add_rounded_card(slide11, px, py, Inches(5.8), Inches(1.7), VERY_LIGHT_BLUE)
    add_accent_bar(slide11, px, py, Inches(0.08), Inches(1.7), PRIMARY_BLUE)
    add_textbox(slide11, px + Inches(0.3), py + Inches(0.15), Inches(5.3), Inches(0.4),
                label, 20, BLACK, True, font_name=BODY_FONT)
    add_textbox(slide11, px + Inches(0.3), py + Inches(0.7), Inches(5.3), Inches(0.8),
                desc, 16, MED_GRAY, False, font_name=BODY_FONT, line_spacing=1.3)


# ════════════════════════════════════════════════════════════════
# SLIDE 12 - Thank You / Q&A
# ════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide12, PRIMARY_BLUE)
add_shape_bg(slide12, Inches(0), Inches(0), Inches(4.5), SLIDE_HEIGHT, CYAN_ACCENT)

add_textbox(slide12, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
            "感谢聆听", 52, WHITE, True, PP_ALIGN.CENTER, font_name=TITLE_FONT)

add_accent_bar(slide12, Inches(5.2), Inches(3.3), Inches(3), Inches(0.06), CYAN_ACCENT)

add_textbox(slide12, Inches(1.2), Inches(3.5), Inches(10), Inches(0.8),
            "智慧校园 -- 让管理更简单，让服务更温暖", 24, RGBColor(0xBB, 0xDE, 0xFB),
            False, PP_ALIGN.CENTER, font_name=BODY_FONT)

add_textbox(slide12, Inches(1.2), Inches(5.2), Inches(10), Inches(0.6),
            "Q & A", 28, RGBColor(0x93, 0xC5, 0xFD), True, PP_ALIGN.CENTER, font_name=BODY_FONT)

add_accent_bar(slide12, Inches(0), Inches(7.2), SLIDE_WIDTH, Inches(0.08), CYAN_ACCENT)


# Save
output_path = r"D:\编程\智慧校园\presentation.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
